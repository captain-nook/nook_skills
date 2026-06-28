# -*- coding: utf-8 -*-
from __future__ import annotations
import argparse, json, re, zipfile
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor

BLACKS = {"151510", "000000", "171713"}

def rgb_hex(color):
    try:
        rgb = color.rgb
        if rgb is None: return None
        return str(rgb).upper()
    except Exception:
        return None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--map")
    args = ap.parse_args()
    pptx = Path(args.pptx)
    cmap = {}
    if args.map and Path(args.map).exists():
        data=json.loads(Path(args.map).read_text(encoding='utf-8-sig'))
        cmap={int(x['page_no']): x for x in data}
    prs=Presentation(pptx)
    rows=[]
    issues=[]
    for idx, slide in enumerate(prs.slides, 1):
        text_boxes=0; chars=0; pics=0; shapes=0; black_fills=0; long_texts=[]
        for shp in slide.shapes:
            shapes += 1
            if getattr(shp, 'has_text_frame', False):
                txt='\n'.join(p.text for p in shp.text_frame.paragraphs).strip()
                if txt:
                    text_boxes += 1
                    chars += len(txt)
                    if len(txt) > 34:
                        long_texts.append(txt[:40])
            if shp.shape_type == 13:
                pics += 1
            try:
                if shp.fill and shp.fill.fore_color:
                    hx=rgb_hex(shp.fill.fore_color)
                    if hx in BLACKS:
                        black_fills += 1
            except Exception:
                pass
        comp=cmap.get(idx,{}).get('component_id','')
        variant=cmap.get(idx,{}).get('variant','')
        page_issues=[]
        if black_fills >= 4: page_issues.append(f"black_blocks={black_fills}")
        if text_boxes >= 14: page_issues.append(f"many_text_boxes={text_boxes}")
        if chars >= 120: page_issues.append(f"dense_chars={chars}")
        if chars <= 38 and pics <= 1 and idx not in [1,20,28]: page_issues.append(f"too_empty_chars={chars}")
        if shapes >= 45: page_issues.append(f"many_shapes={shapes}")
        if long_texts: page_issues.append("long_text="+" | ".join(long_texts[:2]))
        rows.append({"page":idx,"component":comp,"variant":variant,"text_boxes":text_boxes,"chars":chars,"pics":pics,"shapes":shapes,"black_fills":black_fills,"issues":page_issues})
        if page_issues: issues.append(rows[-1])
    print(json.dumps({"total_pages":len(prs.slides),"issue_pages":issues,"rows":rows}, ensure_ascii=False, indent=2))

if __name__=='__main__': main()
