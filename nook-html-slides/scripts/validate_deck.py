# -*- coding: utf-8 -*-
from __future__ import annotations

import argparse
import re
import zipfile
from pathlib import Path
from pptx import Presentation

FORBIDDEN = ["真实截图后补","这里保留证据槽","不伪造","后续替换","REAL SCREENSHOT ONLY","component","slot_id","文件路径","构建","占位"]


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--expected-pages", type=int)
    args = ap.parse_args()
    pptx = Path(args.pptx)
    prs = Presentation(pptx)
    slide_count = len(prs.slides)
    if args.expected_pages and slide_count != args.expected_pages:
        raise SystemExit(f"FAIL slide_count {slide_count} != {args.expected_pages}")
    with zipfile.ZipFile(pptx, "r") as z:
        names = z.namelist()
        slide_xml = [n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n)]
        media = [n for n in names if n.startswith("ppt/media/")]
        hits = {}
        for n in slide_xml:
            txt = z.read(n).decode("utf-8", "ignore")
            bad = [w for w in FORBIDDEN if w in txt]
            if bad:
                hits[n] = bad
        if hits:
            raise SystemExit(f"FAIL forbidden text: {hits}")
    density = []
    for idx, slide in enumerate(prs.slides, 1):
        chars = 0
        pics = 0
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False):
                chars += len("".join(p.text for p in shp.text_frame.paragraphs))
            if shp.shape_type == 13:
                pics += 1
        if chars > 420:
            density.append((idx, chars))
    if density:
        raise SystemExit(f"FAIL text density: {density}")
    print(f"OK slides={slide_count} xml_slides={len(slide_xml)} media={len(media)} forbidden_hits=0")


if __name__ == "__main__":
    main()
