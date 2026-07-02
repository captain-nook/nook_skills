from __future__ import annotations

import json
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SKELETON_PATH = ROOT / "global-component-skeleton-v1.json"

THEMES = {
    "bright-street-dance": {
        "out": ROOT / "bright-street-dance" / "component-library-v0.1.0.pptx",
        "bg": "F4FBFF",
        "panel": "FFFFFF",
        "soft": "E8F6FF",
        "ink": "101014",
        "muted": "4F5963",
        "a": "21A8FF",
        "b": "FF4FA3",
        "c": "B8F13C",
        "d": "FFD33D",
        "dark": "101014",
        "white": "FFFFFF",
        "title_font": "Arial Black",
        "body_font": "Microsoft YaHei UI",
        "num_font": "Arial Black",
        "mono_font": "Consolas",
        "radius": True,
        "line_w": 3.0,
        "shadow_scale": 1.25,
        "code": "STICKER",
    },
    "oriental-dark-yaji": {
        "out": ROOT / "oriental-dark-yaji" / "component-library-v0.1.0.pptx",
        "bg": "17110C",
        "panel": "241A13",
        "soft": "2F241A",
        "ink": "F4E8CF",
        "muted": "BFAE92",
        "a": "B33A2E",
        "b": "C69A55",
        "c": "6E7F68",
        "d": "8A5E3B",
        "dark": "0E0A07",
        "white": "F8F1E2",
        "title_font": "SimSun",
        "body_font": "FangSong",
        "num_font": "Georgia",
        "mono_font": "Consolas",
        "radius": False,
        "line_w": 1.2,
        "shadow_scale": 0.45,
        "code": "YAJI",
    },
}

ALIGN = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def rgb(hexstr: str) -> RGBColor:
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


class SkeletonReplayRenderer:
    def __init__(self, theme: dict, skeleton: dict):
        self.t = theme
        self.skeleton = skeleton
        self.prs = Presentation()
        self.prs.slide_width = Inches(skeleton["page_size"]["width_in"])
        self.prs.slide_height = Inches(skeleton["page_size"]["height_in"])

    def token(self, value: str | None) -> str:
        if not value:
            return self.t["ink"]
        return self.t.get(value, value)

    def font(self, value: str | None) -> str:
        if not value:
            return self.t["body_font"]
        return self.t.get(value, value)

    def align(self, value: str | None):
        if not value:
            return PP_ALIGN.LEFT
        key = str(value).lower()
        if key.startswith("center"):
            return PP_ALIGN.CENTER
        if key.startswith("right"):
            return PP_ALIGN.RIGHT
        return ALIGN.get(key, PP_ALIGN.LEFT)

    def render(self) -> None:
        for component in self.skeleton["components"]:
            slide = self.blank_slide()
            for op in component["ops"]:
                self.dispatch(slide, op)

    def blank_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(self.t["bg"])
        return slide

    def dispatch(self, slide, op: dict) -> None:
        name = op["op"]
        if name == "page_header":
            self.page_header(slide, op["code"], op["title"], op["subtitle"], op["folio"])
        elif name == "text_box":
            self.text_box(
                slide,
                op["value"],
                op["x"],
                op["y"],
                op["w"],
                op["h"],
                font=op.get("font"),
                size=op.get("size", 18),
                color=op.get("color"),
                bold=op.get("bold", False),
                align=op.get("align"),
                line_spacing=op.get("line_spacing"),
            )
        elif name == "bullet_text":
            self.bullet_text(slide, op["lines"], op["x"], op["y"], op["w"], op["h"], op.get("size", 15), op.get("color"))
        elif name == "shadow_rect":
            self.shadow_rect(slide, op["x"], op["y"], op["w"], op["h"], op.get("color"), op.get("dx", 0.06), op.get("dy", 0.06))
        elif name == "rect":
            self.rect(slide, op["x"], op["y"], op["w"], op["h"], op.get("fill"), op.get("line"), op.get("lw", 2.0), op.get("radius", False), op.get("dash", False))
        elif name == "line":
            self.line(slide, op["x1"], op["y1"], op["x2"], op["y2"], op.get("color"), op.get("width", 2.0), op.get("dash", False))
        elif name == "image_slot":
            self.image_slot(slide, op["x"], op["y"], op["w"], op["h"], op["code"], op.get("caption", ""), op.get("accent"), op.get("fill"))
        elif name == "big_num":
            self.big_num(slide, op["value"], op["x"], op["y"], op.get("size", 66))
        elif name == "tiny_rule":
            self.tiny_rule(slide, op["x"], op["y"], op["w"], op.get("color"))
        elif name == "label":
            self.label(slide, op["value"], op["x"], op["y"], op.get("w", 1.7))
        elif name == "component_note":
            self.component_note(slide, op["value"])
        else:
            raise ValueError(f"unknown skeleton op: {name}")

    def text_box(
        self,
        slide,
        value,
        x,
        y,
        w,
        h,
        font=None,
        size=18,
        color=None,
        bold=False,
        align=None,
        valign=MSO_ANCHOR.TOP,
        line_spacing=None,
    ):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        tf.vertical_anchor = valign
        for idx, part in enumerate(str(value).split("\n")):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = part
            p.alignment = self.align(align)
            if line_spacing:
                p.line_spacing = line_spacing
            for run in p.runs:
                run.font.name = self.font(font)
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = rgb(self.token(color))
        return tb

    def bullet_text(self, slide, lines, x, y, w, h, size=15, color=None):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        for i, value in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "· " + value
            p.line_spacing = 0.9
            for run in p.runs:
                run.font.name = self.font("body_font")
                run.font.size = Pt(size)
                run.font.color.rgb = rgb(self.token(color or "muted"))
        return tb

    def rect(self, slide, x, y, w, h, fill="panel", line="ink", line_w=2.0, radius=False, dash=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if self.t["radius"] or radius else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(self.token(fill))
        shp.line.color.rgb = rgb(self.token(line or "ink"))
        shp.line.width = Pt(line_w)
        if dash:
            shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        return shp

    def shadow_rect(self, slide, x, y, w, h, color="b", dx=0.06, dy=0.06):
        scale = self.t.get("shadow_scale", 1.0)
        return self.rect(slide, x + dx * scale, y + dy * scale, w, h, color, color, 0.1, False, False)

    def line(self, slide, x1, y1, x2, y2, color="ink", width=2.0, dash=False):
        shp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        shp.line.color.rgb = rgb(self.token(color))
        shp.line.width = Pt(width)
        if dash:
            shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        return shp

    def page_header(self, slide, code, title, subtitle, page):
        code = code.replace("GIANT TITLE", self.t["code"])
        self.text_box(slide, code, 0.72, 0.30, 5.8, 0.24, font="mono_font", size=9.5, color="muted")
        self.text_box(slide, title, 0.70, 0.62, 7.6, 0.76, font="title_font", size=31, color="ink", bold=True)
        self.text_box(slide, subtitle, 10.0, 0.42, 5.25, 0.44, font="mono_font", size=8.5, color="muted", align="right")
        self.text_box(slide, page, 0.72, 8.36, 4.2, 0.18, font="mono_font", size=7.8, color="muted")
        self.rect(slide, 14.08, 8.18, 1.18, 0.32, "panel", "muted", 1.0, False, True)
        self.text_box(slide, "LOGO SLOT", 14.16, 8.25, 1.02, 0.10, font="mono_font", size=6.3, color="muted", align="center")

    def label(self, slide, value, x, y, w=1.7):
        self.rect(slide, x, y, w, 0.28, "ink", "ink", 0.5, False, False)
        self.text_box(slide, value, x + 0.06, y + 0.075, w - 0.12, 0.10, font="mono_font", size=6.6, color="bg", align="center")

    def image_slot(self, slide, x, y, w, h, code, caption="", accent="b", fill="panel"):
        self.shadow_rect(slide, x, y, w, h, accent, 0.06, 0.06)
        self.rect(slide, x, y, w, h, fill or "panel", "ink", 2.0, False, False)
        self.rect(slide, x + 0.22, y + 0.32, w - 0.44, h - 0.72, "bg", "muted", 1.2, False, True)
        self.label(slide, code, x + 0.20, y + 0.16, 1.55)
        if caption:
            self.text_box(slide, caption, x + 0.28, y + h - 0.32, w - 0.56, 0.12, font="mono_font", size=6.5, color="muted", align="right")

    def big_num(self, slide, value, x, y, size=66):
        self.text_box(slide, value, x + 0.06, y + 0.07, 2.1, 0.75, font="num_font", size=size, color="b", bold=True)
        self.text_box(slide, value, x, y, 2.1, 0.75, font="num_font", size=size, color="ink", bold=True)

    def tiny_rule(self, slide, x, y, w, color="a"):
        return self.rect(slide, x, y, w, 0.05, color, color, 0.1, False, False)

    def component_note(self, slide, value):
        self.rect(slide, 11.55, 7.35, 3.7, 0.62, "panel", "ink", 1.4, False, False)
        self.text_box(slide, value, 11.75, 7.50, 3.3, 0.26, font="body_font", size=10.2, color="muted", align="center")

    def save(self) -> None:
        out = self.t["out"]
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(out)
        with zipfile.ZipFile(out, "r") as zf:
            slides = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
            media = [n for n in zf.namelist() if n.startswith("ppt/media/")]
        print(f"saved={out}")
        print(f"slides={len(slides)} media={len(media)} bytes={out.stat().st_size}")


def main() -> None:
    skeleton = json.loads(SKELETON_PATH.read_text(encoding="utf-8"))
    if skeleton.get("source_mode") != "extracted_from_duotone_component_library_source":
        raise SystemExit("global skeleton must be extracted from duotone source")
    for theme in THEMES.values():
        renderer = SkeletonReplayRenderer(theme, skeleton)
        renderer.render()
        renderer.save()


if __name__ == "__main__":
    main()
