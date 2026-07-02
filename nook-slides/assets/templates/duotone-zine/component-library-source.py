
from pathlib import Path
import zipfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "component-library-v0.4.0.pptx"

PAPER = "F3ECD8"
PAPER_SOFT = "F8F0DD"
INK = "171713"
MUTED = "5E5A51"
BLUE = "3153B8"
BLUE_DEEP = "243D86"
PINK = "E76186"
YELLOW = "E8B13F"
WHITE = "FFFFFF"

TITLE_FONT = "思源宋体 CN Heavy"
BODY_FONT = "霞鹜文楷"
NUM_FONT = "Bodoni 72"

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def rgb(hexstr):
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))

def blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(PAPER)
    return slide

def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)

def set_line(shape, color=INK, width=2.0, dash=None):
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if dash:
        shape.line.dash_style = dash

def no_line(shape):
    shape.line.fill.background()

def rect(slide, x, y, w, h, fill_color=PAPER_SOFT, line_color=INK, line_w=2.0, radius=False, dash=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shp, fill_color)
    set_line(shp, line_color, line_w, dash)
    return shp

def line(slide, x1, y1, x2, y2, color=INK, width=2.0, dash=None):
    shp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shp.line.color.rgb = rgb(color)
    shp.line.width = Pt(width)
    if dash:
        shp.line.dash_style = dash
    return shp

def text_box(slide, text, x, y, w, h, font=BODY_FONT, size=18, color=INK, bold=False, align=None, valign=MSO_ANCHOR.TOP, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    parts = str(text).split("\n")
    for idx, part in enumerate(parts):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = part
        if align:
            p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return tb

def bullet_text(slide, lines, x, y, w, h, size=15, color=MUTED):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "· " + txt
        p.line_spacing = 0.9
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(size)
            run.font.color.rgb = rgb(color)
    return tb

def shadow_rect(slide, x, y, w, h, color=PINK, dx=0.06, dy=0.06):
    shp = rect(slide, x+dx, y+dy, w, h, fill_color=color, line_color=color, line_w=0.1)
    return shp

def page_header(slide, code, title, subtitle, page):
    text_box(slide, code, 0.72, 0.30, 5.8, 0.24, font="Consolas", size=9.5, color=MUTED)
    text_box(slide, title, 0.70, 0.62, 7.6, 0.76, font=TITLE_FONT, size=31, color=INK, bold=True)
    text_box(slide, subtitle, 10.0, 0.42, 5.25, 0.44, font="Consolas", size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)
    text_box(slide, page, 0.72, 8.36, 4.2, 0.18, font="Consolas", size=7.8, color=MUTED)
    logo = rect(slide, 14.08, 8.18, 1.18, 0.32, fill_color=PAPER_SOFT, line_color=MUTED, line_w=1.0, dash=MSO_LINE_DASH_STYLE.DASH)
    text_box(slide, "LOGO SLOT", 14.16, 8.25, 1.02, 0.10, font="Consolas", size=6.3, color=MUTED, align=PP_ALIGN.CENTER)

def label(slide, txt, x, y, w=1.7):
    rect(slide, x, y, w, 0.28, fill_color=INK, line_color=INK, line_w=0.5)
    text_box(slide, txt, x+0.06, y+0.075, w-0.12, 0.10, font="Consolas", size=6.6, color=PAPER, align=PP_ALIGN.CENTER)

def image_slot(slide, x, y, w, h, code, caption="", accent=PINK, fill_color=PAPER_SOFT):
    shadow_rect(slide, x, y, w, h, accent, 0.06, 0.06)
    rect(slide, x, y, w, h, fill_color=fill_color, line_color=INK, line_w=2.0)
    ph = rect(slide, x+0.22, y+0.32, w-0.44, h-0.72, fill_color=PAPER, line_color=MUTED, line_w=1.2, dash=MSO_LINE_DASH_STYLE.DASH)
    label(slide, code, x+0.20, y+0.16, 1.55)
    if caption:
        text_box(slide, caption, x+0.28, y+h-0.32, w-0.56, 0.12, font="Consolas", size=6.5, color=MUTED, align=PP_ALIGN.RIGHT)
    return ph

def big_num(slide, txt, x, y, size=66):
    text_box(slide, txt, x+0.06, y+0.07, 2.1, 0.75, font=NUM_FONT, size=size, color=PINK, bold=True)
    text_box(slide, txt, x, y, 2.1, 0.75, font=NUM_FONT, size=size, color=INK, bold=True)

def tiny_rule(slide, x, y, w, color=BLUE):
    r = rect(slide, x, y, w, 0.05, fill_color=color, line_color=color, line_w=0.1)
    return r

def component_note(slide, text):
    rect(slide, 11.55, 7.35, 3.7, 0.62, fill_color=PAPER_SOFT, line_color=INK, line_w=1.4)
    text_box(slide, text, 11.75, 7.50, 3.3, 0.26, font=BODY_FONT, size=10.2, color=MUTED, align=PP_ALIGN.CENTER)

def slide_01_giant_title():
    slide = blank_slide()
    page_header(slide, "C01 / GIANT TITLE", "巨型标题页", "少字 / 大字 / 强入口", "P01 / GIANT TITLE")
    image_slot(slide, 9.65, 1.55, 4.85, 5.85, "IMG-A", "竖版主图槽 4:5", accent=BLUE)
    text_box(slide, "可编辑\nPPT\n不是截图", 0.90, 1.62, 7.9, 2.95, font=TITLE_FONT, size=52, color=INK, bold=True, line_spacing=0.82)
    tiny_rule(slide, 0.95, 5.05, 1.1, PINK)
    text_box(slide, "它是一套能被拆开、移动、重排、再编辑的出版组件。", 0.92, 5.35, 6.2, 0.46, font=BODY_FONT, size=20, color=MUTED)
    component_note(slide, "适用：封面 / 章节页 / 大观点开场")

def slide_02_one_sentence():
    slide = blank_slide()
    page_header(slide, "C02 / ONE SENTENCE", "一句话判断页", "一句话占主导，不写小作文", "P02 / ONE SENTENCE")
    shadow_rect(slide, 1.0, 1.75, 13.3, 3.35, BLUE, 0.08, 0.08)
    rect(slide, 1.0, 1.75, 13.3, 3.35, fill_color=BLUE_DEEP, line_color=INK, line_w=2.4)
    text_box(slide, "模型不能临场\n发明页面", 1.45, 2.20, 7.2, 1.35, font=TITLE_FONT, size=43, color=WHITE, bold=True, line_spacing=0.86)
    text_box(slide, "它只能从组件白名单里选择、组合、替换内容。", 1.48, 3.95, 8.7, 0.35, font=BODY_FONT, size=20, color=WHITE)
    label(slide, "CORE JUDGEMENT", 11.15, 2.18, 2.2)
    text_box(slide, "深底白字\n最多两行标题", 11.05, 3.0, 2.4, 0.56, font="Consolas", size=11, color=WHITE, align=PP_ALIGN.CENTER)
    component_note(slide, "适用：结论页 / 转场页 / 警示页")

def slide_03_big_number():
    slide = blank_slide()
    page_header(slide, "C03 / BIG NUMBER", "大数字页", "数字是主角，文字只解释", "P03 / BIG NUMBER")
    big_num(slide, "12", 0.9, 1.65, 96)
    text_box(slide, "类组件", 3.55, 2.1, 2.2, 0.48, font=TITLE_FONT, size=28, color=INK, bold=True)
    text_box(slide, "第一套组件库不追求多，追求少而稳定。", 0.98, 4.0, 5.7, 0.44, font=BODY_FONT, size=20, color=MUTED)
    # three metric columns
    for i, (num, title, desc, color) in enumerate([
        ("05", "图片槽", "多图页面的核心", BLUE),
        ("04", "文本框", "少字大字的基础", PINK),
        ("03", "结构图", "时间轴、循环、层级", YELLOW),
    ]):
        x = 7.0 + i*2.65
        shadow_rect(slide, x, 2.0, 2.2, 2.9, color, 0.055, 0.055)
        rect(slide, x, 2.0, 2.2, 2.9, fill_color=PAPER_SOFT, line_color=INK, line_w=2.0)
        text_box(slide, num, x+0.22, 2.28, 1.2, 0.52, font=NUM_FONT, size=34, color=INK, bold=True)
        tiny_rule(slide, x+0.25, 3.05, 0.9, color)
        text_box(slide, title, x+0.25, 3.38, 1.6, 0.25, font=TITLE_FONT, size=18, color=INK, bold=True)
        text_box(slide, desc, x+0.25, 3.82, 1.6, 0.35, font=BODY_FONT, size=12.5, color=MUTED)
    component_note(slide, "适用：数据页 / 进度页 / 章节摘要")

def slide_04_process_cards():
    slide = blank_slide()
    page_header(slide, "C04 / PROCESS CARDS", "流程卡组", "可复制的步骤文本框", "P04 / PROCESS CARDS")
    cards = [
        ("01", "输入", "收集材料", ["文本", "图片", "品牌资产"], BLUE),
        ("02", "处理", "拆成组件", ["内容压缩", "图片入槽", "字体检查"], PINK),
        ("03", "生成", "拼成页面", ["选择版式", "替换文案", "输出 PPT"], YELLOW),
        ("04", "检查", "人工审美", ["可编辑", "不出框", "风格统一"], INK),
    ]
    for i, (num, label_txt, title, bullets, color) in enumerate(cards):
        x = 0.85 + i*3.65
        shadow_rect(slide, x, 2.05, 3.05, 4.25, color, 0.06, 0.06)
        rect(slide, x, 2.05, 3.05, 4.25, fill_color=PAPER_SOFT, line_color=INK, line_w=2.1)
        rect(slide, x, 2.05, 3.05, 0.72, fill_color=color, line_color=INK, line_w=2.1)
        text_box(slide, num, x+0.18, 2.18, 0.72, 0.4, font=NUM_FONT, size=28, color=WHITE if color != YELLOW else INK, bold=True)
        text_box(slide, label_txt.upper(), x+1.05, 2.34, 1.65, 0.14, font="Consolas", size=8.4, color=WHITE if color != YELLOW else INK)
        text_box(slide, title, x+0.32, 3.15, 2.2, 0.32, font=TITLE_FONT, size=21, color=INK, bold=True)
        tiny_rule(slide, x+0.33, 3.70, 1.1, color)
        bullet_text(slide, bullets, x+0.35, 4.10, 2.25, 1.0, size=13.5, color=MUTED)
    component_note(slide, "适用：线性流程 / 方法论步骤 / 操作说明")

def slide_05_timeline():
    slide = blank_slide()
    page_header(slide, "C05 / TIMELINE", "时间轴", "横向主轴 + 事件卡片", "P05 / TIMELINE")
    line(slide, 1.1, 4.45, 14.2, 4.45, color=BLUE, width=9)
    line(slide, 1.1, 4.58, 14.2, 4.58, color=PINK, width=4)
    events = [
        ("1920s", 1.05, 2.0, BLUE, "源头", "小众出版\n手作表达"),
        ("1976", 5.0, 1.55, YELLOW, "浪潮", "ZINE 文化\n开始扩散"),
        ("1991", 8.95, 4.92, PINK, "转折", "社群、音乐\n与女性主义"),
        ("2010s+", 12.25, 1.85, BLUE, "复兴", "书展、馆藏\n与独立出版"),
    ]
    for year, x, y, color, title, desc in events:
        big_num(slide, year, x, y-0.85, 34 if len(year)>4 else 46)
        shadow_rect(slide, x+0.1, y+0.35, 2.5, 1.25, color, 0.05, 0.05)
        rect(slide, x+0.1, y+0.35, 2.5, 1.25, fill_color=INK, line_color=INK, line_w=1.8)
        rect(slide, x+0.1, y+0.35, 2.5, 0.38, fill_color=color, line_color=INK, line_w=1.6)
        text_box(slide, title, x+0.25, y+0.43, 2.15, 0.14, font=TITLE_FONT, size=13, color=INK if color == YELLOW else WHITE, align=PP_ALIGN.CENTER)
        text_box(slide, desc, x+0.28, y+0.88, 2.0, 0.38, font=BODY_FONT, size=11, color=WHITE, align=PP_ALIGN.CENTER)
        line(slide, x+1.35, y+1.60, x+1.35, 4.43, color=INK, width=2.2)
    component_note(slide, "适用：历史线 / 项目里程碑 / 版本演进")

def slide_06_cycle():
    slide = blank_slide()
    page_header(slide, "C06 / CYCLE", "循环图", "3-4 步循环，不做复杂飞轮", "P06 / CYCLE")
    cx, cy = 8.0, 4.35
    nodes = [
        ("01", "输入", 8.0, 1.85, BLUE),
        ("02", "生成", 11.2, 4.35, PINK),
        ("03", "检查", 8.0, 6.75, YELLOW),
        ("04", "迭代", 4.8, 4.35, INK),
    ]
    # connector arrows approximation: thick lines + triangles
    line(slide, 8.9, 2.35, 10.45, 3.72, color=BLUE, width=3)
    line(slide, 10.7, 5.0, 8.9, 6.28, color=PINK, width=3)
    line(slide, 7.05, 6.3, 5.5, 5.0, color=YELLOW, width=3)
    line(slide, 5.55, 3.7, 7.05, 2.35, color=INK, width=3)
    for no, txt, x, y, color in nodes:
        shadow_rect(slide, x-0.75, y-0.45, 1.5, 0.9, color, 0.05, 0.05)
        rect(slide, x-0.75, y-0.45, 1.5, 0.9, fill_color=PAPER_SOFT, line_color=INK, line_w=2.0)
        text_box(slide, no, x-0.55, y-0.25, 0.45, 0.22, font=NUM_FONT, size=18, color=color if color != INK else INK, bold=True)
        text_box(slide, txt, x-0.05, y-0.18, 0.55, 0.18, font=TITLE_FONT, size=14, color=INK, bold=True)
    text_box(slide, "一轮不求完美\n只求能被检查", 0.95, 2.25, 4.2, 0.88, font=TITLE_FONT, size=30, color=INK, bold=True, line_spacing=0.9)
    text_box(slide, "循环图用于表达反复校准的工作流，不用于承载大量文字。", 1.0, 4.1, 4.3, 0.42, font=BODY_FONT, size=17, color=MUTED)
    component_note(slide, "适用：反馈循环 / 生成-检查-迭代 / 闭环系统")

def slide_07_hierarchy():
    slide = blank_slide()
    page_header(slide, "C07 / HIERARCHY", "层级图", "堆叠卡 / 上下游 / 金字塔替代", "P07 / HIERARCHY")
    levels = [
        ("视觉体系", "颜色、字体、纹理、母版", BLUE, 2.0, 1.65, 12.0),
        ("组件白名单", "文本框、图片槽、结构图", PINK, 2.65, 3.00, 10.7),
        ("页面版式", "封面、目录、流程、时间线", YELLOW, 3.30, 4.35, 9.4),
        ("具体内容", "替换文字与素材，不临场发明", INK, 3.95, 5.70, 8.1),
    ]
    for title, desc, color, x, y, w in levels:
        shadow_rect(slide, x, y, w, 0.82, color, 0.06, 0.06)
        rect(slide, x, y, w, 0.82, fill_color=PAPER_SOFT, line_color=INK, line_w=2.1)
        rect(slide, x, y, 0.30, 0.82, fill_color=color, line_color=INK, line_w=1.0)
        text_box(slide, title, x+0.55, y+0.18, 2.2, 0.18, font=TITLE_FONT, size=17, color=INK, bold=True)
        text_box(slide, desc, x+3.0, y+0.24, w-3.3, 0.15, font=BODY_FONT, size=13.5, color=MUTED)
    component_note(slide, "适用：系统分层 / 依赖关系 / 优先级")

def slide_08_comparison():
    slide = blank_slide()
    page_header(slide, "C08 / COMPARISON", "对比图", "左右对照 / 三栏对照", "P08 / COMPARISON")
    items = [
        ("随机生成", "风格漂移\n文字容易满\n图片乱摆", PINK, "不采用"),
        ("组件驱动", "组件白名单\n图片先入槽\n输出可检查", BLUE, "采用"),
    ]
    for i, (title, body, color, tag) in enumerate(items):
        x = 1.05 + i*7.05
        shadow_rect(slide, x, 1.85, 5.95, 4.75, color, 0.07, 0.07)
        rect(slide, x, 1.85, 5.95, 4.75, fill_color=PAPER_SOFT, line_color=INK, line_w=2.4)
        rect(slide, x, 1.85, 5.95, 0.68, fill_color=color, line_color=INK, line_w=2.0)
        text_box(slide, title, x+0.45, 2.85, 4.8, 0.48, font=TITLE_FONT, size=30, color=INK, bold=True)
        text_box(slide, body, x+0.50, 3.80, 4.0, 0.9, font=BODY_FONT, size=21, color=MUTED, line_spacing=0.9)
        label(slide, tag.upper(), x+0.45, 5.60, 1.7)
    component_note(slide, "适用：方案对比 / 前后对比 / 取舍判断")

def slide_09_image_text():
    slide = blank_slide()
    page_header(slide, "C09 / IMAGE + TEXT", "图文混排页", "图大，字少，文字只做锚点", "P09 / IMAGE TEXT")
    image_slot(slide, 0.95, 1.55, 8.1, 5.95, "IMG-B", "横版主图 / 16:9 或近似", accent=BLUE)
    text_box(slide, "图片先决定页面气质", 9.65, 1.90, 4.65, 0.8, font=TITLE_FONT, size=33, color=INK, bold=True, line_spacing=0.9)
    tiny_rule(slide, 9.72, 3.05, 1.0, PINK)
    text_box(slide, "这类页面只允许一段短说明。图片不是装饰，它是主内容。", 9.72, 3.40, 4.1, 0.65, font=BODY_FONT, size=18, color=MUTED)
    label(slide, "CAPTION", 9.72, 5.55, 1.25)
    text_box(slide, "图片来源 / 生成提示 / 风格化说明", 9.72, 6.00, 4.0, 0.28, font="Consolas", size=9, color=MUTED)
    component_note(slide, "适用：多图章节 / 案例页 / 插画解释页")

def slide_10_collage():
    slide = blank_slide()
    page_header(slide, "C10 / COLLAGE", "图片拼贴页", "主图 + 辅图，不平均分配", "P10 / COLLAGE")
    image_slot(slide, 0.95, 1.55, 6.6, 5.85, "IMG-A", "主图槽", accent=BLUE)
    image_slot(slide, 8.05, 1.55, 3.1, 2.45, "IMG-E1", "辅图", accent=PINK)
    image_slot(slide, 11.55, 1.55, 3.1, 2.45, "IMG-E2", "辅图", accent=YELLOW)
    image_slot(slide, 8.05, 4.45, 6.6, 2.95, "IMG-B", "宽幅辅图", accent=PINK)
    text_box(slide, "拼贴不是网格平均分。主图必须明确，辅图只做证据。", 8.12, 7.62, 6.2, 0.25, font=BODY_FONT, size=14, color=MUTED)
    component_note(slide, "适用：案例合集 / 视觉证据 / 多图讲述")

def slide_11_master_ip_logo():
    slide = blank_slide()
    page_header(slide, "C11 / MASTER ASSETS", "IP / Logo 母版页", "Logo 原样，IP 风格化后入槽", "P11 / MASTER ASSETS")
    shadow_rect(slide, 0.95, 1.65, 4.55, 2.55, BLUE, 0.06, 0.06)
    rect(slide, 0.95, 1.65, 4.55, 2.55, fill_color=PAPER_SOFT, line_color=INK, line_w=2.1)
    text_box(slide, "Logo：原样角标", 1.25, 2.0, 2.8, 0.32, font=TITLE_FONT, size=22, color=INK, bold=True)
    rect(slide, 3.88, 2.18, 1.12, 0.36, fill_color=PAPER, line_color=MUTED, line_w=1.1, dash=MSO_LINE_DASH_STYLE.DASH)
    text_box(slide, "USER LOGO", 3.98, 2.30, 0.92, 0.08, font="Consolas", size=6.2, color=MUTED, align=PP_ALIGN.CENTER)
    text_box(slide, "不重绘、不改色、不套颗粒。\n只做小尺寸固定角标。", 1.25, 2.72, 3.5, 0.55, font=BODY_FONT, size=14, color=MUTED)
    shadow_rect(slide, 6.15, 1.65, 8.45, 2.55, PINK, 0.06, 0.06)
    rect(slide, 6.15, 1.65, 8.45, 2.55, fill_color=BLUE_DEEP, line_color=INK, line_w=2.1)
    text_box(slide, "IP：必须先风格化", 6.48, 2.0, 3.0, 0.32, font=TITLE_FONT, size=22, color=WHITE, bold=True)
    text_box(slide, "用户提供原图 → 确认保留特征 → 图生图转成双色老印刷 → 抠图入槽", 6.48, 2.65, 6.7, 0.36, font=BODY_FONT, size=14, color=WHITE)
    image_slot(slide, 0.95, 4.75, 3.6, 2.25, "IP-A", "主视觉 IP 槽", accent=BLUE)
    image_slot(slide, 4.95, 4.75, 3.0, 2.25, "IP-B", "旁白 IP 槽", accent=PINK)
    image_slot(slide, 8.35, 4.75, 2.7, 2.25, "IP-C", "动作 IP 槽", accent=YELLOW)
    image_slot(slide, 11.45, 4.75, 3.1, 2.25, "IP-D", "流程小图标槽", accent=INK)
    component_note(slide, "适用：品牌母版 / 角色资产 / 统一出图规则")

def slide_12_image_generation_assets():
    slide = blank_slide()
    page_header(slide, "C12 / IMAGE GENERATION", "出图资产页", "把生图变成可控页面元素", "P12 / IMAGE ASSETS")
    text_box(slide, "出图不是自由发挥\n而是进入槽位", 0.95, 1.65, 6.2, 1.35, font=TITLE_FONT, size=37, color=INK, bold=True, line_spacing=0.86)
    text_box(slide, "每张图都要先确定：用途、比例、是否抠图、是否需要图生图统一风格。", 1.0, 3.38, 5.8, 0.46, font=BODY_FONT, size=18, color=MUTED)
    rules = [
        ("A", "原图输入", "用户图 / 素材图 / IP"),
        ("B", "风格化", "双色、纸感、网点、错位"),
        ("C", "抠图裁切", "透明 PNG / 固定比例"),
        ("D", "进入槽位", "A/B/C/D/E 图片容器"),
    ]
    for i, (n, title, desc) in enumerate(rules):
        x = 7.35 + (i%2)*3.45
        y = 1.75 + (i//2)*2.45
        color = [BLUE, PINK, YELLOW, INK][i]
        shadow_rect(slide, x, y, 2.95, 1.65, color, 0.055, 0.055)
        rect(slide, x, y, 2.95, 1.65, fill_color=PAPER_SOFT, line_color=INK, line_w=2.0)
        rect(slide, x+0.22, y+0.22, 0.52, 0.52, fill_color=color, line_color=INK, line_w=1.3)
        text_box(slide, n, x+0.31, y+0.34, 0.34, 0.15, font=NUM_FONT, size=13, color=WHITE if color != YELLOW else INK, align=PP_ALIGN.CENTER)
        text_box(slide, title, x+0.92, y+0.30, 1.6, 0.22, font=TITLE_FONT, size=16, color=INK, bold=True)
        text_box(slide, desc, x+0.30, y+0.92, 2.35, 0.32, font=BODY_FONT, size=12, color=MUTED)
    shadow_rect(slide, 1.0, 5.35, 5.85, 1.5, PINK, 0.06, 0.06)
    rect(slide, 1.0, 5.35, 5.85, 1.5, fill_color=BLUE_DEEP, line_color=INK, line_w=2.0)
    text_box(slide, "运行时必须问：有没有 logo？有没有 IP？是否允许风格化和抠图？", 1.35, 5.85, 5.1, 0.30, font=BODY_FONT, size=15.5, color=WHITE, align=PP_ALIGN.CENTER)
    component_note(slide, "适用：出图任务 / 图生图 / 图片资产收口")

def slide_13_card_group():
    slide = blank_slide()
    page_header(slide, "C13 / CARD GROUP", "并列卡片组", "同字段 / 同权重 / 无方向关系", "P13 / CARD GROUP")
    text_box(slide, "卡片只表达并列", 0.92, 1.52, 7.2, 0.72, font=TITLE_FONT, size=34, color=INK, bold=True)
    text_box(slide, "没有箭头、主轴、闭环或父子关系。", 0.96, 2.38, 7.2, 0.36, font=BODY_FONT, size=18, color=MUTED)
    entries = [
        ("01", "内容", "锁定页数和可见文字", BLUE),
        ("02", "形式", "确认页面任务和变体", PINK),
        ("03", "资产", "确认截图、插图和数据", YELLOW),
    ]
    for i, (num, title, body, color) in enumerate(entries):
        x = 1.05 + i * 4.78
        shadow_rect(slide, x, 3.18, 3.95, 3.0, color, 0.065, 0.065)
        rect(slide, x, 3.18, 3.95, 3.0, fill_color=PAPER_SOFT, line_color=INK, line_w=2.0)
        text_box(slide, num, x + 0.28, 3.54, 0.92, 0.46, font=NUM_FONT, size=30, color=INK, bold=True)
        tiny_rule(slide, x + 0.30, 4.25, 1.05, color)
        text_box(slide, title, x + 0.32, 4.64, 2.9, 0.42, font=TITLE_FONT, size=22, color=INK, bold=True)
        text_box(slide, body, x + 0.34, 5.34, 3.1, 0.38, font=BODY_FONT, size=14.5, color=MUTED)
    component_note(slide, "适用：并列模块 / 案例 / 角色 / 交付物")

slides = [
    slide_01_giant_title,
    slide_02_one_sentence,
    slide_03_big_number,
    slide_04_process_cards,
    slide_05_timeline,
    slide_06_cycle,
    slide_07_hierarchy,
    slide_08_comparison,
    slide_09_image_text,
    slide_10_collage,
    slide_11_master_ip_logo,
    slide_12_image_generation_assets,
    slide_13_card_group,
]

for fn in slides:
    fn()

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)

with zipfile.ZipFile(OUT, "r") as zf:
    names = zf.namelist()
    slide_xml = [n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
    media = [n for n in names if n.startswith("ppt/media/")]
    print(f"saved={OUT}")
    print(f"bytes={OUT.stat().st_size}")
    print(f"slides={len(slide_xml)}")
    print(f"media_files={len(media)}")
    for n in slide_xml:
        data = zf.read(n).decode("utf-8", "ignore")
        print(n, "shapes", data.count("<p:sp>"), "pics", data.count("<p:pic>"))

