# -*- coding: utf-8 -*-
"""Template-driven editable PPTX builder for nook-slides.

Input: deck_plan.json
Output: editable PPTX + build log + internal component map

This script is intentionally boring: it does not invent layouts. It maps a locked
plan to maintained components and prints progress for the user.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
import time
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from plan_contract import validate_plan_contract

FORBIDDEN_VISIBLE = [
    "真实截图后补", "这里保留证据槽", "不伪造", "后续替换",
    "REAL SCREENSHOT ONLY", "component", "slot_id", "文件路径", "构建", "占位",
]

C = {
    "paper": "F4EDDA", "paper2": "EFE6CF", "ink": "151510", "muted": "625F55",
    "blue": "3156B8", "pink": "E75D82", "gold": "E3AD3D", "black": "151510",
    "cream": "F8F1DE", "line": "1B1A16", "white": "FFFFFF"
}
TITLE_FONT = "思源宋体 CN Heavy"
BODY_FONT = "霞鹜文楷"
MONO_FONT = "Courier New"
NUM_FONT = "Bodoni 72"

WIDE_W = Inches(16)
WIDE_H = Inches(9)


def now() -> str:
    return time.strftime("%H:%M:%S")


def progress(stage: str, message: str) -> None:
    print(f"[{now()}] {stage}: {message}", flush=True)


def rgb(hexstr: str) -> RGBColor:
    value = C.get(hexstr, hexstr).strip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:], 16))


def color_name(idx: int) -> str:
    return ["blue", "pink", "gold", "black", "cream"][idx % 5]

def shadow_name(accent: str) -> str:
    # Keep offset-shadow visibly different from the front panel fill/strip.
    return {"blue": "pink", "pink": "blue", "gold": "blue", "black": "gold", "cream": "pink"}.get(accent, "pink")



def clean_text(text, keep_breaks: bool = False) -> str:
    """Return locked visible copy without rewriting characters or punctuation."""
    return str(text or "")

@dataclass
class BuildContext:
    prs: Presentation
    out_dir: Path
    crop_dir: Path
    logo: Optional[Path]
    blank: Any
    component_map: List[Dict[str, Any]]


class ZineRenderer:
    def __init__(self, ctx: BuildContext):
        self.ctx = ctx

    # ---------- primitives ----------
    def add_text(self, slide, text, x, y, w, h, size=34, font=BODY_FONT, color="ink", bold=False,
                 align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, spacing=1.0, keep_breaks=False):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
        tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = valign
        text = clean_text(text, keep_breaks)
        for i, line in enumerate(str(text or "").split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = spacing
            r = p.add_run(); r.text = line
            r.font.name = font
            r.font.bold = bold
            r.font.size = Pt(size)
            r.font.color.rgb = rgb(color)
        return box

    def rect(self, slide, x, y, w, h, fill="cream", line="line", lw=2.2):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
        shp.line.color.rgb = rgb(line); shp.line.width = Pt(lw)
        return shp

    def line(self, slide, x1, y1, x2, y2, color="blue", width=3):
        shp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        shp.line.color.rgb = rgb(color); shp.line.width = Pt(width)
        return shp

    def bg(self, slide, page_no: int, total: int):
        self.rect(slide, 0, 0, 16, 9, "paper", "paper", 0)
        self.add_text(slide, f"{page_no:02d} / {total}", 0.35, 8.55, 1.0, 0.18, 7.5, MONO_FONT, "muted")
        if self.ctx.logo and self.ctx.logo.exists():
            try:
                slide.shapes.add_picture(str(self.ctx.logo), Inches(15.15), Inches(8.25), width=Inches(0.48))
            except Exception as exc:
                progress("WARN", f"logo insert failed on p{page_no}: {exc}")

    def panel(self, slide, x, y, w, h, accent="blue", fill="cream", lw=2.2):
        self.rect(slide, x + 0.08, y + 0.10, w, h, shadow_name(accent), shadow_name(accent), 0)
        return self.rect(slide, x, y, w, h, fill, "line", lw)

    def crop_cover(self, img_path: Path, ratio: float, name: str, focus_x: float = 0.5, focus_y: float = 0.5) -> Path:
        im = Image.open(img_path).convert("RGB")
        iw, ih = im.size
        current = iw / ih
        if current > ratio:
            nw = int(ih * ratio)
            left = int(max(0, min(iw - nw, iw * focus_x - nw / 2)))
            box = (left, 0, left + nw, ih)
        else:
            nh = int(iw / ratio)
            top = int(max(0, min(ih - nh, ih * focus_y - nh / 2)))
            box = (0, top, iw, top + nh)
        out = self.ctx.crop_dir / f"{name}.jpg"
        im.crop(box).save(out, quality=92)
        return out

    def add_picture_fit(self, slide, img, x, y, w, h, fit="cover", name="img", focus_x=0.5, focus_y=0.5):
        img = Path(img)
        if not img.exists():
            return False
        if fit == "cover":
            out = self.crop_cover(img, w / h, name, focus_x, focus_y)
            slide.shapes.add_picture(str(out), Inches(x), Inches(y), Inches(w), Inches(h))
        else:
            im = Image.open(img)
            iw, ih = im.size
            scale = min(w / (iw / 100), h / (ih / 100))
            pw, ph = (iw / 100) * scale, (ih / 100) * scale
            slide.shapes.add_picture(str(img), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))
        return True

    def image_slot(self, slide, x, y, w, h, label, img=None, accent="pink", fit="cover", name="slot", focus_x=0.5):
        self.panel(slide, x, y, w, h, accent, "paper2", 2.0)
        if img and Path(img).exists():
            self.add_picture_fit(slide, img, x + 0.08, y + 0.08, w - 0.16, h - 0.16, fit, name, focus_x)
        else:
            self.add_text(slide, label or "待替换真实截图", x + 0.25, y + h / 2 - 0.18, w - 0.5, 0.36, 20, BODY_FONT, "muted", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    def chip(self, slide, text, x, y, w, h, fill="black", color="white", size=22):
        self.rect(slide, x, y, w, h, fill, "line", 1.6)
        self.add_text(slide, text, x + 0.08, y + 0.08, w - 0.16, h - 0.16, size, BODY_FONT, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    def card(self, slide, title, body="", x=1, y=1, w=3, h=2, accent="blue", dark=False):
        """Reusable zine card with a real text-box feel, not a UI button."""
        self.panel(slide, x, y, w, h, accent, "cream", 2.0)
        # Top strip restores the reference text-box grammar, but stays slim.
        strip_h = min(0.34, max(0.18, h * 0.16))
        strip_fill = "gold" if accent == "cream" else accent
        self.rect(slide, x, y, w, strip_h, strip_fill, strip_fill, 0)
        title_size = 20 if w >= 2.6 else 17 if len(str(title)) <= 8 else 15
        head_color = "white" if strip_fill in ["blue", "pink", "black"] else "ink"
        self.add_text(slide, title, x + 0.24, y + 0.06, w - 0.48, strip_h - 0.06, title_size, TITLE_FONT, head_color, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        if body:
            body_len = len(str(body))
            body_size = 20 if body_len <= 14 and h >= 1.6 else 17 if body_len <= 24 else 14 if body_len <= 40 else 12
            self.add_text(slide, body, x + 0.30, y + strip_h + 0.30, w - 0.60, max(0.28, h - strip_h - 0.45), body_size, BODY_FONT, "muted", False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, 0.9)

    def new_slide(self, page: Dict[str, Any], total: int):
        slide = self.ctx.prs.slides.add_slide(self.ctx.blank)
        self.bg(slide, int(page["page_no"]), total)
        self.ctx.component_map.append({
            "page_no": page["page_no"],
            "component_id": page.get("component"),
            "variant": page.get("variant", ""),
            "asset_status": [slot.get("status", "") for slot in page.get("image_slots", [])],
        })
        return slide

    # ---------- component renderers ----------
    def render(self, page: Dict[str, Any], total: int):
        component = page.get("component", "")
        method = {
            "zine.hero-title": self.hero_title,
            "zine.statement": self.statement,
            "zine.big-number": self.big_number,
            "zine.process-cards": self.process_cards,
            "zine.card-group": self.card_group,
            "zine.timeline": self.timeline,
            "zine.cycle": self.cycle,
            "zine.hierarchy": self.hierarchy,
            "zine.comparison": self.comparison,
            "zine.image-text": self.image_text,
            "zine.image-collage": self.image_collage,
            "zine.brand-master": self.brand_master,
            "zine.generated-asset": self.generated_asset,
        }.get(component)
        if not method:
            raise ValueError(f"Unsupported component: {component} on page {page.get('page_no')}")
        progress("RENDER", f"p{page['page_no']:02d} {component} / {page.get('title','')}")
        return method(page, total)

    def first_image(self, page: Dict[str, Any]) -> Dict[str, Any]:
        slots = page.get("image_slots") or []
        return slots[0] if slots else {}

    def hero_title(self, page, total):
        s = self.new_slide(page, total)
        slot = self.first_image(page)
        title = page.get("title", "")
        subtitle = page.get("subtitle") or page.get("body", "")
        variant = page.get("variant", "")
        if slot:
            if variant == "case-opener":
                self.add_text(s, title, 1.05, 1.05, 6.00, 1.35, 33, TITLE_FONT, "ink", True, spacing=0.86)
                self.image_slot(s, 8.35, 0.95, 6.20, 6.15, slot.get("placeholder", ""), slot.get("path"), "gold", slot.get("fit", "cover"), f"p{page['page_no']}", slot.get("focus_x", 0.5))
                chips = page.get("chips", [])[:6]
                for i, t in enumerate(chips):
                    row, col_i = divmod(i, 2)
                    x = 1.15 + col_i * 2.65
                    y = 4.25 + row * 0.78
                    col = color_name(i)
                    self.chip(s, t, x, y, 2.28, 0.58, col, "white" if col in ["blue", "pink", "black"] else "ink", 18)
                return
            self.add_text(s, title, 1.25, 1.10, 5.90, 2.95, 50, TITLE_FONT, "ink", True, spacing=0.84, keep_breaks=(int(page.get("page_no", 0)) == 1))
            if subtitle:
                self.add_text(s, subtitle, 1.32, 4.35, 5.35, 0.55, 22, BODY_FONT, "muted")
            self.image_slot(s, 8.05, 1.05, 6.45, 6.25, slot.get("placeholder", ""), slot.get("path"), "blue", slot.get("fit", "cover"), f"p{page['page_no']}", slot.get("focus_x", 0.5))
        else:
            self.add_text(s, title, 2.05, 1.35, 11.9, 2.65, 54, TITLE_FONT, "ink", True, PP_ALIGN.CENTER, spacing=0.86)
            if subtitle:
                self.add_text(s, subtitle, 3.3, 4.45, 9.4, 0.55, 24, BODY_FONT, "muted", align=PP_ALIGN.CENTER)
        chips = page.get("chips", [])[:6]
        if chips:
            total_w = min(9.8, len(chips)*1.42 + (len(chips)-1)*0.32)
            x0 = (16-total_w)/2
            for i, t in enumerate(chips):
                self.chip(s, t, x0 + i*1.74, 6.55, 1.28, 0.46, color_name(i), "white" if color_name(i) in ["blue", "pink", "black"] else "ink", 15)
    def statement(self, page, total):
        s = self.new_slide(page, total)
        title = page.get("title", "")
        body = page.get("body", "")
        chips = page.get("chips", [])[:5]
        if chips:
            self.add_text(s, title, 1.35, 1.40, 6.70, 1.05, 36, TITLE_FONT, "ink", True, spacing=0.88)
            if body:
                self.add_text(s, body, 1.42, 2.78, 6.10, 0.58, 21, BODY_FONT, "muted")
            self.panel(s, 8.38, 2.35, 5.20, 4.15, "blue", "cream", 2.0)
            for i, t in enumerate(chips):
                col = color_name(i)
                y = 2.78 + i * 0.64
                self.rect(s, 8.78, y, 0.16, 0.34, col, col, 0)
                self.add_text(s, t, 9.12, y - 0.02, 3.70, 0.30, 18, BODY_FONT, "ink", True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        else:
            self.add_text(s, title, 2.10, 1.45, 11.80, 2.20, 50, TITLE_FONT, "ink", True, PP_ALIGN.CENTER, spacing=0.88)
            if body:
                self.add_text(s, body, 3.00, 4.05, 10.00, 0.68, 25, BODY_FONT, "muted", align=PP_ALIGN.CENTER)

    def big_number(self, page, total):
        s = self.new_slide(page, total)
        number = str(page.get("number", ""))
        self.add_text(s, number, 0.95, 1.25, 4.5, 1.2, 88, NUM_FONT, "ink", True)
        self.add_text(s, page.get("title", ""), 0.98, 3.15, 5.4, 0.7, 32, TITLE_FONT, "ink", True)
        if page.get("body"):
            self.add_text(s, page["body"], 1.0, 4.20, 5.8, 0.62, 21, BODY_FONT, "muted")
        cards = page.get("cards", [])[:3]
        for i,c in enumerate(cards):
            x = 7.0 + i*2.65
            col = c.get("accent") or color_name(i)
            self.panel(s, x, 2.0, 2.2, 2.9, col, "cream", 2.0)
            self.add_text(s, f"{i+1:02d}", x+0.22, 2.28, 1.2, 0.52, 34, NUM_FONT, "ink", True)
            self._tiny_rule(s, x+0.25, 3.05, 0.9, col)
            self.add_text(s, c.get("title", ""), x+0.25, 3.38, 1.6, 0.25, 17, TITLE_FONT, "ink", True)
            if c.get("body"):
                self.add_text(s, c.get("body", ""), x+0.25, 3.82, 1.6, 0.35, 12.5, BODY_FONT, "muted")

    def process_cards(self, page, total):
        s = self.new_slide(page, total)
        self.add_text(s, page.get("title", ""), 0.90, 0.78, 11.8, 0.92, 38, TITLE_FONT, "ink", True)
        if page.get("body"):
            self.add_text(s, page["body"], 0.95, 1.78, 10.8, 0.40, 20, BODY_FONT, "muted")
        cards = page.get("cards", [])[:5]
        if len(cards) <= 3:
            self._template_three_cards(s, cards)
        elif len(cards) in [4,5]:
            self._template_process_cards(s, cards)
        else:
            raise ValueError(f"page {page.get('page_no')}: process-cards supports up to 5 visible cards")
        if page.get("footer"):
            self.add_text(s, page["footer"], 2.4, 7.72, 11.2, 0.34, 20, BODY_FONT, "muted", align=PP_ALIGN.CENTER)

    def card_group(self, page, total):
        """Render parallel peers only: no arrows, axis, loop, or parent-child claim."""
        s = self.new_slide(page, total)
        self.add_text(s, page.get("title", ""), 0.90, 0.78, 11.8, 0.92, 38, TITLE_FONT, "ink", True)
        if page.get("body"):
            self.add_text(s, page["body"], 0.95, 1.78, 11.8, 0.40, 20, BODY_FONT, "muted")
        cards = page.get("cards", [])
        n = len(cards)
        width = 3.85 if n == 3 else 2.82
        gap = 0.82 if n == 3 else 0.58
        total_w = n * width + (n - 1) * gap
        x0 = (16 - total_w) / 2
        for i, card in enumerate(cards):
            x = x0 + i * (width + gap)
            accent = card.get("accent") or color_name(i)
            self.panel(s, x, 2.55, width, 3.65, accent, "cream", 2.0)
            self.add_text(s, f"{i + 1:02d}", x + 0.28, 2.88, 0.72, 0.42, 28, NUM_FONT, "ink", True)
            self._tiny_rule(s, x + 0.30, 3.58, 1.05, accent)
            self.add_text(s, card.get("title", ""), x + 0.32, 3.95, width - 0.64, 0.45, 20, TITLE_FONT, "ink", True)
            if card.get("body"):
                self.add_text(s, card["body"], x + 0.34, 4.65, width - 0.68, 0.78, 15, BODY_FONT, "muted")

    def _tiny_rule(self, slide, x, y, w, color="blue"):
        return self.rect(slide, x, y, w, 0.055, color, color, 0)

    def _template_three_cards(self, slide, cards):
        width = 3.85
        gap = 0.85
        total_w = len(cards) * width + max(0, len(cards)-1) * gap
        x0 = (16 - total_w) / 2
        for i, c in enumerate(cards):
            x = x0 + i * (width + gap)
            col = c.get("accent") or color_name(i)
            self.panel(slide, x, 2.35, width, 4.15, col, "cream", 2.1)
            self.rect(slide, x, 2.35, width, 0.68, col, col, 0)
            self.add_text(slide, f"{i+1:02d}", x+0.24, 2.46, 0.72, 0.36, 26, NUM_FONT, "white" if col in ["blue","pink","black"] else "ink", True)
            self.add_text(slide, c.get("title", ""), x+0.38, 3.38, width-0.76, 0.50, 22, TITLE_FONT, "ink", True)
            self._tiny_rule(slide, x+0.40, 4.05, 1.05, col)
            if c.get("body"):
                self.add_text(slide, c.get("body", ""), x+0.42, 4.42, width-0.84, 1.15, 16, BODY_FONT, "muted", False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 0.9)

    def _template_process_cards(self, slide, cards):
        n = len(cards)
        width = 3.05 if n == 4 else 2.55
        gap = 0.60 if n == 4 else 0.42
        total_w = n * width + (n-1) * gap
        x0 = (16 - total_w) / 2
        for i, c in enumerate(cards):
            x = x0 + i * (width + gap)
            col = c.get("accent") or color_name(i)
            self.panel(slide, x, 2.55, width, 4.45, col, "cream", 2.1)
            strip_col = "gold" if col == "cream" else col
            self.rect(slide, x, 2.55, width, 0.72, strip_col, strip_col, 0)
            self.add_text(slide, f"{i+1:02d}", x+0.18, 2.68, 0.72, 0.40, 28, NUM_FONT, "white" if strip_col in ["blue","pink","black"] else "ink", True)
            self.add_text(slide, c.get("label", "STEP").upper(), x+1.05, 2.84, max(0.8,width-1.25), 0.16, 8.4, MONO_FONT, "white" if strip_col in ["blue","pink","black"] else "ink")
            self.add_text(slide, c.get("title", ""), x+0.32, 3.65, width-0.64, 0.42, 20 if n == 4 else 17, TITLE_FONT, "ink", True)
            self._tiny_rule(slide, x+0.33, 4.20, 1.1, col)
            if c.get("body"):
                self.add_text(slide, c.get("body", ""), x+0.35, 4.60, width-0.70, 1.05, 13.5 if n == 5 else 14.5, BODY_FONT, "muted", False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 0.9)

    def arrow(self, slide, x, y, w, h, fill="black"):
        shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
        shp.line.color.rgb = rgb(fill); shp.line.width = Pt(0)
        return shp
    def timeline(self, page, total):
        s = self.new_slide(page, total)
        self.add_text(s, page.get("title", ""), 1.05, 0.62, 13.8, 0.82, 34, TITLE_FONT, "ink", True)
        if page.get("body"):
            self.add_text(s, page["body"], 1.12, 1.58, 10.8, 0.38, 19, BODY_FONT, "muted")
        cards = page.get("cards", [])[:5]
        if len(cards) == 3:
            self._timeline_three_cpr(s, cards)
            return
        self.line(s, 1.10, 4.50, 13.65, 4.50, "blue", 9)
        self.line(s, 1.10, 4.63, 13.65, 4.63, "pink", 4)
        self.arrow(s, 14.35, 4.37, 0.46, 0.36, "blue")
        if len(cards) == 5:
            self._timeline_five_layers(s, cards)
            return
        xs = [1.05, 5.00, 8.95, 12.25][:len(cards)]
        ys = [2.10, 4.85, 2.10, 4.85][:len(cards)]
        for i, c in enumerate(cards):
            x, y = xs[i], ys[i]
            col = c.get("accent") or color_name(i)
            year = c.get("year") or c.get("title", str(i+1))
            self.panel(s, x+0.10, y+0.35, 2.35, 1.25, col, "cream", 1.8)
            self.rect(s, x+0.10, y+0.35, 0.22, 1.25, col, col, 0)
            self.add_text(s, str(year), x+0.42, y+0.52, 1.85, 0.26, 16, TITLE_FONT, "ink", True)
            if c.get("body"):
                self.add_text(s, c.get("body", ""), x+0.42, y+0.98, 1.80, 0.28, 12.5, BODY_FONT, "muted", False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, 0.9)
            self.line(s, x+1.28, y+1.60, x+1.28, 4.50, "ink", 1.8)

    def _timeline_three_cpr(self, slide, cards):
        specs = [(1.15, 2.75, 3.70, 2.75), (6.15, 2.75, 3.70, 2.75), (11.15, 2.75, 3.70, 2.75)]
        for i, c in enumerate(cards):
            x, y, w, h = specs[i]
            col = c.get("accent") or color_name(i)
            self.panel(slide, x, y, w, h, col, "cream", 2.1)
            self.rect(slide, x, y, 0.52, h, col, col, 0)
            self.add_text(slide, c.get("title", ""), x + 0.85, y + 0.48, 0.78, 0.58, 42, NUM_FONT, "ink", True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
            if c.get("body"):
                self.add_text(slide, c.get("body", ""), x + 0.88, y + 1.42, w - 1.25, 0.65, 15.5, BODY_FONT, "muted", False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, 0.9)
            if i < len(cards) - 1:
                self.arrow(slide, x + w + 0.35, y + 1.20, 0.52, 0.34, "black")

    def _timeline_five_layers(self, slide, cards):
        top_y, bottom_y = 2.65, 5.55
        specs = [
            (0.95, top_y, 2.42, 1.22),
            (3.95, bottom_y, 2.42, 1.22),
            (6.78, top_y, 2.42, 1.22),
            (9.60, bottom_y, 2.42, 1.22),
            (12.45, top_y, 2.42, 1.22),
        ]
        for i, c in enumerate(cards):
            x, y, w, h = specs[i]
            col = c.get("accent") or color_name(i)
            node_x = x + w / 2
            node_y = 4.56
            if y < node_y:
                self.line(slide, node_x, y + h, node_x, node_y, "ink", 1.6)
            else:
                self.line(slide, node_x, node_y, node_x, y, "ink", 1.6)
            self.rect(slide, node_x - 0.13, node_y - 0.13, 0.26, 0.26, col, "line", 1.2)
            self.panel(slide, x, y, w, h, col, "cream", 1.9)
            self.rect(slide, x, y, 0.18, h, col, col, 0)
            self.add_text(slide, f"{i+1:02d}", x + 0.34, y + 0.22, 0.45, 0.20, 14, NUM_FONT, col if col != "black" else "ink", True)
            self.add_text(slide, c.get("title", ""), x + 0.84, y + 0.20, w - 1.10, 0.24, 16, TITLE_FONT, "ink", True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
            if c.get("body"):
                self.add_text(slide, c.get("body", ""), x + 0.34, y + 0.68, w - 0.60, 0.25, 12.5, BODY_FONT, "muted", False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, 0.9)
    def cycle(self, page, total):
        s = self.new_slide(page, total)
        self.add_text(s, page.get("title", ""), 1.30, 0.72, 13.40, 0.82, 38, TITLE_FONT, "ink", True, PP_ALIGN.CENTER)
        if page.get("body"):
            self.add_text(s, page["body"], 1.95, 1.76, 12.1, 0.42, 19, BODY_FONT, "muted", align=PP_ALIGN.CENTER)
        nodes = page.get("cards", [])[:4]
        positions = [(8.0,2.72),(11.2,4.85),(8.0,6.75),(4.8,4.85)]
        connectors = [(8.9,3.22,10.45,4.22,"blue"),(10.7,5.5,8.9,6.28,"pink"),(7.05,6.3,5.5,5.5,"gold"),(5.55,4.2,7.05,3.22,"black")]
        for x1,y1,x2,y2,col in connectors:
            self.line(s,x1,y1,x2,y2,col,3)
        for i,c in enumerate(nodes):
            x,y = positions[i]
            col = c.get("accent") or color_name(i)
            self.panel(s, x-0.75, y-0.45, 1.5, 0.9, col, "cream", 2.0)
            self.add_text(s, f"{i+1:02d}", x-0.55, y-0.25, 0.45, 0.22, 18, NUM_FONT, col if col != "black" else "ink", True)
            self.add_text(s, c.get("title", ""), x-0.05, y-0.18, 0.62, 0.18, 13, TITLE_FONT, "ink", True)
        self.add_text(s, page.get("footer", ""), 4.00, 7.62, 8.00, 0.38, 18, BODY_FONT, "muted", align=PP_ALIGN.CENTER)

    def hierarchy(self, page, total):
        s = self.new_slide(page, total)
        self.add_text(s, page.get("title", ""), 0.85, 0.72, 12.6, 0.88, 37, TITLE_FONT, "ink", True)
        if page.get("body"):
            self.add_text(s, page["body"], 0.92, 1.72, 11.2, 0.38, 20, BODY_FONT, "muted")
        cards = page.get("cards", [])[:5]
        self._template_layer_stack(s, cards)
        if page.get("footer"):
            self.add_text(s, page["footer"], 2.35, 7.55, 11.0, 0.34, 20, BODY_FONT, "muted", align=PP_ALIGN.CENTER)

    def _template_layer_stack(self, slide, cards):
        n = len(cards)
        if n <= 4:
            specs = [(1.45,2.75,13.1),(2.10,3.88,11.8),(2.75,5.01,10.5),(3.40,6.14,9.2)]
            h = 1.02
        else:
            specs = [(1.25,2.68,13.5),(1.80,3.70,12.4),(2.35,4.72,11.3),(2.90,5.74,10.2),(3.45,6.76,9.1)]
            h = 0.92
        for i,c in enumerate(cards):
            x,y,w = specs[i]
            col = c.get("accent") or color_name(i)
            self.panel(slide, x, y, w, h, col, "cream", 2.1)
            strip_col = "gold" if col == "cream" else col
            self.rect(slide, x, y, 0.36, h, strip_col, strip_col, 0)
            self.add_text(slide, c.get("title", ""), x+0.62, y+0.22, 2.30, 0.24, 17.5, TITLE_FONT, "ink", True)
            if c.get("body"):
                self.add_text(slide, c.get("body", ""), x+3.05, y+0.25, max(1.0,w-3.45), 0.24, 14.2, BODY_FONT, "muted")

    def comparison(self, page, total):
        s = self.new_slide(page, total)
        self.add_text(s, page.get("title", ""), 0.85, 0.72, 12.6, 0.88, 37, TITLE_FONT, "ink", True)
        if page.get("body"):
            self.add_text(s, page["body"], 0.92, 1.72, 11.2, 0.38, 20, BODY_FONT, "muted")
        cards = page.get("cards", [])[:3]
        if len(cards) <= 2:
            width, gap = 5.95, 1.10
        else:
            width, gap = 4.20, 0.65
        total_w = len(cards)*width + max(0,len(cards)-1)*gap
        x0 = (16-total_w)/2
        for i,c in enumerate(cards):
            x = x0 + i*(width+gap)
            col = c.get("accent") or color_name(i)
            self.panel(s, x, 2.15, width, 4.75, col, "cream", 2.4)
            self.rect(s, x, 2.15, width, 0.68, col, col, 0)
            self.add_text(s, c.get("title", ""), x+0.45, 3.05, width-0.90, 0.48, 26 if len(cards)<=2 else 21, TITLE_FONT, "ink", True)
            if c.get("body"):
                self.add_text(s, c.get("body", ""), x+0.50, 3.95, width-1.0, 1.05, 18 if len(cards)<=2 else 15.5, BODY_FONT, "muted", False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 0.9)
            label = c.get("tag") or c.get("label") or ("采用" if i == len(cards)-1 else "")
            if label:
                self.add_text(s, label.upper(), x+0.45, 5.82, 1.7, 0.20, 10, MONO_FONT, col, True)
        if page.get("footer"):
            self.add_text(s, page["footer"], 2.35, 7.55, 11.0, 0.34, 20, BODY_FONT, "muted", align=PP_ALIGN.CENTER)

    def image_text(self, page, total):
        s = self.new_slide(page, total)
        slot = self.first_image(page)
        has_real_image = bool(slot.get("path")) and Path(slot.get("path")).exists()
        # Always use the sample's large image + right text geometry. Missing screenshots remain inside the image slot.
        self.image_slot(s, 0.95, 1.55, 8.10, 5.95, slot.get("placeholder", "待替换真实截图"), slot.get("path") if has_real_image else None, slot.get("accent", "blue"), slot.get("fit", "contain"), f"p{page['page_no']}", slot.get("focus_x", 0.5))
        self.add_text(s, page.get("title", ""), 9.65, 1.90, 4.65, 0.72, 27, TITLE_FONT, "ink", True, spacing=0.9)
        self._tiny_rule(s, 9.72, 3.12, 1.0, slot.get("accent", "pink"))
        if page.get("body"):
            self.add_text(s, page["body"], 9.72, 3.45, 4.10, 0.95, 17.5, BODY_FONT, "muted", False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 0.9)
        chips = page.get("chips", [])[:4]
        for i,t in enumerate(chips):
            self.add_text(s, str(t), 9.72, 5.18 + i*0.38, 3.9, 0.22, 12.5, MONO_FONT, "muted")

    def image_collage(self, page, total):
        s = self.new_slide(page, total)
        self.add_text(s, page.get("title", ""), 1.05, 0.72, 13.8, 0.82, 36, TITLE_FONT, "ink", True)
        if page.get("body"):
            self.add_text(s, page["body"], 1.12, 1.66, 11.8, 0.40, 21, BODY_FONT, "muted")
        slots = page.get("image_slots", [])
        # Centered collage: main visual block spans the page center instead of hugging left.
        coords = [(1.15,2.05,6.55,5.85),(8.20,2.05,3.05,2.55),(11.60,2.05,3.05,2.55),(8.20,5.05,6.45,2.85)]
        for i, slot in enumerate(slots[:4]):
            x,y,w,h = coords[i]
            self.image_slot(s, x, y, w, h, slot.get("placeholder", "待替换真实截图"), slot.get("path"), slot.get("accent", color_name(i)), slot.get("fit", "contain"), f"p{page['page_no']}_{i}")

    def brand_master(self, page, total):
        s = self.new_slide(page, total)
        self.add_text(s, page.get("title", "Logo 与 IP 运行规则"), 0.7, 0.62, 8.8, 0.72, 40, TITLE_FONT, "ink", True)
        self.card(s, "Logo：原样角标", "不重绘、不改色、不套颗粒。", 0.9, 2.0, 4.8, 2.0, "blue")
        self.card(s, "IP：先风格化", "用户提供原图后，图生图统一到当前视觉体系。", 6.4, 2.0, 5.7, 2.0, "pink")

    def generated_asset(self, page, total):
        s = self.new_slide(page, total)
        self.add_text(s, page.get("title", "出图资产"), 0.85, 0.72, 10.8, 0.88, 38, TITLE_FONT, "ink", True)
        if page.get("body"):
            self.add_text(s, page["body"], 0.92, 1.55, 10.8, 0.44, 21, BODY_FONT, "muted")
        cards = page.get("cards", [])[:4]
        self._template_process_cards(s, cards if len(cards)>=4 else cards + [{}]*(4-len(cards)))

    def _cards_grid(self, slide, cards: List[Dict[str, Any]], x0, y0, w, h, cols=3, gap=0.8):
        for i, c in enumerate(cards):
            x = x0 + (i % cols) * (w + gap)
            y = y0 + (i // cols) * (h + 0.55)
            col = c.get("accent") or color_name(i)
            self.card(slide, c.get("title", ""), c.get("body", ""), x, y, w, h, col, dark=(col == "black"))


def load_plan(path: Path) -> Dict[str, Any]:
    progress("LOAD", f"reading plan {path}")
    return json.loads(path.read_text(encoding="utf-8-sig"))



def validate_plan(plan: Dict[str, Any], plan_path: Optional[Path] = None) -> None:
    progress("CHECK", "validating locked plan structure")
    validate_plan_contract(plan, plan_path)
    if plan.get("theme") != "duotone-zine":
        raise ValueError(
            f"theme {plan.get('theme')!r} has a shared schema contract, but no approved PPTX renderer yet"
        )
    pages = plan["pages"]
    for p in pages:
        no = p["page_no"]
        visible = json.dumps({k: p.get(k) for k in ["title", "subtitle", "body", "footer", "cards", "chips"]}, ensure_ascii=False)
        for word in FORBIDDEN_VISIBLE:
            if word in visible:
                raise ValueError(f"page {no} contains forbidden construction text: {word}")
    progress("CHECK", f"plan ok: {len(pages)} pages")


def build(plan: Dict[str, Any], out: Path) -> Dict[str, Any]:
    progress("BUILD", "initializing editable PPTX")
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    crop_dir = out.parent / f"_{out.stem}_crops"
    crop_dir.mkdir(parents=True, exist_ok=True)
    logo = Path(plan.get("logo", "")) if plan.get("logo") else None
    ctx = BuildContext(prs=prs, out_dir=out.parent, crop_dir=crop_dir, logo=logo, blank=prs.slide_layouts[6], component_map=[])
    renderer = ZineRenderer(ctx)
    pages = plan["pages"]
    total = len(pages)
    for page in pages:
        renderer.render(page, total)
    progress("SAVE", f"writing {out}")
    prs.save(out)
    return {"component_map": ctx.component_map, "crop_dir": str(crop_dir)}


def validate_pptx(pptx: Path, expected_pages: int) -> Dict[str, Any]:
    progress("VALIDATE", "checking package, slide count, media, forbidden text")
    prs = Presentation(pptx)
    if len(prs.slides) != expected_pages:
        raise ValueError(f"slide count mismatch: {len(prs.slides)} != {expected_pages}")
    hits = {}
    with zipfile.ZipFile(pptx, "r") as z:
        names = z.namelist()
        slides = [n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n)]
        media = [n for n in names if n.startswith("ppt/media/")]
        for n in slides:
            txt = z.read(n).decode("utf-8", errors="ignore")
            bad = [w for w in FORBIDDEN_VISIBLE if w in txt]
            if bad:
                hits[n] = bad
    if hits:
        raise ValueError(f"forbidden construction text found: {hits}")
    # object-level text density check
    density_issues = []
    for idx, slide in enumerate(prs.slides, 1):
        chars = 0
        pics = 0
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False):
                chars += len("".join(p.text for p in shp.text_frame.paragraphs))
            if shp.shape_type == 13:
                pics += 1
        if chars > 420:
            density_issues.append((idx, chars))
    if density_issues:
        raise ValueError(f"text density too high: {density_issues}")
    progress("VALIDATE", f"ok: slides={expected_pages}, media={len(media)}")
    return {"slides": expected_pages, "media": len(media)}


def write_artifacts(plan: Dict[str, Any], out: Path, result: Dict[str, Any], validation: Dict[str, Any]) -> None:
    map_path = out.with_name(out.stem + "-component-map.json")
    log_path = out.with_name(out.stem + "-build-log.md")
    map_path.write_text(json.dumps(result["component_map"], ensure_ascii=False, indent=2), encoding="utf-8")
    log = [
        f"# {out.name} 构建日志",
        "",
        f"输出文件：`{out}`",
        f"主题：{plan.get('theme', '')}",
        f"页数：{validation['slides']}",
        f"媒体文件数：{validation['media']}",
        "",
        "## 可见进度机制",
        "",
        "本文件由 `scripts/build_deck.py` 生成。构建时会在终端输出 LOAD / CHECK / BUILD / RENDER / SAVE / VALIDATE 阶段，用户可以看到当前正在处理哪一页、哪个组件。",
        "",
        "## 规则",
        "",
        "- 文字只来自 deck_plan。",
        "- 页面只映射到维护过的组件。",
        "- 缺失截图只在图片槽内显示极简待替换提示。",
        "- 构建阶段检查施工污染词和文字密度。",
    ]
    log_path.write_text("\n".join(log) + "\n", encoding="utf-8")
    progress("DONE", f"artifacts: {out.name}, {map_path.name}, {log_path.name}")


def main(argv: Optional[List[str]] = None) -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--plan", required=True, help="Path to deck_plan.json")
    ap.add_argument("--out", required=True, help="Output pptx path")
    args = ap.parse_args(argv)
    try:
        plan_path = Path(args.plan)
        out = Path(args.out)
        out.parent.mkdir(parents=True, exist_ok=True)
        plan = load_plan(plan_path)
        validate_plan(plan, plan_path)
        result = build(plan, out)
        validation = validate_pptx(out, len(plan["pages"]))
        write_artifacts(plan, out, result, validation)
        return 0
    except Exception as exc:
        progress("FAILED", str(exc))
        return 1


if __name__ == "__main__":
    raise SystemExit(main())

















